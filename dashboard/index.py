from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new PowerPoint presentation
prs = Presentation()

# Define slide size (16:9 ratio)
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[5])

# Define colors for the quadrants (light pastel shades)
quadrant_colors = [
    RGBColor(198, 224, 252),  # Top left (low impact, high probability)
    RGBColor(255, 204, 188),  # Top right (high impact, high probability)
    RGBColor(219, 238, 202),  # Bottom left (low impact, low probability)
    RGBColor(255, 249, 184),  # Bottom right (high impact, low probability)
]

# Define quadrant positions (2x2 matrix)
slide_margin = Inches(1.5)  # More margin for cleaner spacing
matrix_size = Inches(7)
half_size = matrix_size / 2

# Draw background quadrants
for i, (x_offset, y_offset) in enumerate([(0, 0), (half_size, 0), (0, half_size), (half_size, half_size)]):
    shape = slide.shapes.add_shape(
        1, slide_margin + x_offset, slide_margin + y_offset, half_size, half_size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = quadrant_colors[i]
    shape.line.color.rgb = RGBColor(255, 255, 255)

# Draw axis labels with rotated probability text
textbox = slide.shapes.add_textbox(slide_margin - Inches(1), slide_margin + matrix_size / 2 - Inches(0.5), Inches(1.2), Inches(1))
text_frame = textbox.text_frame
p = text_frame.add_paragraph()
p.text = "Probability"
p.font.bold = True
p.font.size = Pt(18)
textbox.rotation = 90  # Rotate text for vertical axis

textbox = slide.shapes.add_textbox(slide_margin + matrix_size / 2 - Inches(1), slide_margin + matrix_size + Inches(0.3), Inches(3), Inches(1))
text_frame = textbox.text_frame
p = text_frame.add_paragraph()
p.text = "Impact"
p.font.bold = True
p.font.size = Pt(18)

# Draw quadrant separation lines (thin and subtle)
for x, y, width, height in [
    (slide_margin + half_size, slide_margin, Inches(0.05), matrix_size),  # Vertical
    (slide_margin, slide_margin + half_size, matrix_size, Inches(0.05)),  # Horizontal
]:
    line = slide.shapes.add_shape(1, x, y, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(100, 100, 100)

# Risk placement with improved alignment
risks = [
    ("Funding & Budget Volatility", "Sports bonds, tax incentives, micro-sponsorship", (0.3, 0.15)),  # Medium Probability, High Impact
    ("Doping & Reputational Risk", "Frequent tests, WADA alignment, blockchain logs", (0.7, 0.15)),  # High Probability, High Impact
    ("Governance & Bureaucratic Hurdles", "Single-window clearances, transparent dashboards", (0.3, 0.55)),  # Medium Probability, Medium Impact
    ("Socio-Cultural Barriers & Dropouts", "Language campaigns, parental outreach, flexible scheduling", (0.3, 0.85)),  # Medium Probability, Low Impact
    ("Tech Adoption Challenges", "Subsidies for HPC equipment, ongoing tech training", (0.7, 0.85)),  # Low Probability, High Impact
]

for risk, mitigation, (x_rel, y_rel) in risks:
    x = slide_margin + x_rel * matrix_size - Inches(1.2)
    y = slide_margin + y_rel * matrix_size - Inches(0.5)
    textbox = slide.shapes.add_textbox(x, y, Inches(2.8), Inches(1.2))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = risk
    p.font.bold = True
    p.font.size = Pt(14)
    p.space_after = Pt(3)

    p = text_frame.add_paragraph()
    p.text = f"• {mitigation}"
    p.font.size = Pt(12)

# Save presentation
ppt_filename = "Matrix_Only.pptx"
prs.save(ppt_filename)

# Project Note:
# This script creates a visually appealing single slide with a 2x2 Probability-Impact matrix.
# To use in another presentation, open "Matrix_Only.pptx" and copy/paste the matrix elements.
# Ensure you have python-pptx installed: pip install python-pptx

print(f"PowerPoint saved as {ppt_filename}")
